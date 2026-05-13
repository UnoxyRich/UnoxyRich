#!/usr/bin/env python3
"""Generate machine tool market analysis PPT for Q1 2026."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color Palette ──────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # 深蓝底
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 青蓝强调
ACCENT2   = RGBColor(0xFF, 0xA5, 0x00)   # 橙色强调
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
CARD_BG   = RGBColor(0x1A, 0x2E, 0x44)   # 卡片深背景
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT = RGBColor(0xB0, 0xBE, 0xC5)
GREEN     = RGBColor(0x06, 0xD6, 0xA0)
RED       = RGBColor(0xEF, 0x47, 0x6F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ───────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, alpha=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(14), bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, font_size=Pt(13), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False,
             level=0, bullet_char=None):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.level = level
    if bullet_char:
        p.text = ""
    run = p.add_run()
    run.text = (bullet_char + " " if bullet_char else "") + text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def bg_dark(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=DARK_BG)


def bg_light(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT_BG)


def accent_bar(slide, w=Inches(0.07), color=ACCENT):
    add_rect(slide, Inches(0.45), Inches(1.05), w, Inches(0.55), fill_color=color)


def section_title(slide, title, subtitle=None, light=False):
    tc = TEXT_DARK if light else WHITE
    gc = RGBColor(0x55, 0x66, 0x77) if light else GRAY_TEXT
    accent_bar(slide)
    add_textbox(slide, title,
                Inches(0.65), Inches(1.0), Inches(11.5), Inches(0.65),
                font_size=Pt(26), bold=True, color=tc)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.65), Inches(1.65), Inches(11.5), Inches(0.4),
                    font_size=Pt(14), color=gc, italic=True)


def divider_line(slide, y, light=False):
    color = RGBColor(0x30, 0x50, 0x70) if not light else RGBColor(0xCC, 0xD8, 0xE8)
    add_rect(slide, Inches(0.45), y, Inches(12.4), Inches(0.015), fill_color=color)


# ── Card helper ────────────────────────────────────────────────────────────────

def card(slide, l, t, w, h, bg=CARD_BG, radius=False):
    r = add_rect(slide, l, t, w, h, fill_color=bg)
    return r


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)

# Top accent strip
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
# Bottom accent strip
add_rect(sl, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=ACCENT2)

# Left vertical bar
add_rect(sl, Inches(0.45), Inches(1.2), Inches(0.12), Inches(3.8), fill_color=ACCENT)

# Title block
add_textbox(sl, "爆单·缺货·涨价",
            Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.1),
            font_size=Pt(46), bold=True, color=ACCENT)

add_textbox(sl, "三大驱动因素深度总结",
            Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.7),
            font_size=Pt(32), bold=True, color=WHITE)

add_textbox(sl, "为何 2026 年一季度机床市场异常火热？",
            Inches(0.8), Inches(2.95), Inches(11.5), Inches(0.55),
            font_size=Pt(22), bold=False, color=GRAY_TEXT)

divider_line(sl, Inches(3.65))

# Meta info cards
for i, (label, val) in enumerate([
    ("报告日期", "2026年5月13日"),
    ("报告性质", "行业深度分析"),
    ("核心关键词", "AI赋能 · 共振复苏 · 国产替代"),
]):
    x = Inches(0.8 + i * 4.1)
    card(sl, x, Inches(3.8), Inches(3.9), Inches(0.9), bg=CARD_BG)
    add_textbox(sl, label, x + Inches(0.15), Inches(3.82),
                Inches(3.6), Inches(0.3), font_size=Pt(10), color=GRAY_TEXT)
    add_textbox(sl, val, x + Inches(0.15), Inches(4.08),
                Inches(3.6), Inches(0.55), font_size=Pt(14), bold=True, color=ACCENT)

# Core viewpoint box
add_rect(sl, Inches(0.45), Inches(4.9), Inches(12.4), Inches(1.8),
         fill_color=RGBColor(0x05, 0x3A, 0x5A), line_color=ACCENT, line_width=Pt(1.2))
add_textbox(sl, "【核心观点】",
            Inches(0.7), Inches(4.95), Inches(11.8), Inches(0.4),
            font_size=Pt(13), bold=True, color=ACCENT)
add_textbox(sl,
    "2026年一季度，中国机床行业结束自2023年以来的筑底期，迎来景气度明显拐点。"
    "市场呈现【爆单、缺货、涨价】罕见火热局面，核心驱动力为：AI产业链强力拉动、"
    "多下游需求共振复苏，以及供应链短缺催化的国产替代加速。",
    Inches(0.7), Inches(5.35), Inches(11.8), Inches(1.2),
    font_size=Pt(14), color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – 目录
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)

add_textbox(sl, "目 录", Inches(0.45), Inches(0.15), Inches(12), Inches(0.7),
            font_size=Pt(18), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

items = [
    ("01", "市场现状",     "爆单·缺货·涨价三重信号详解",   ACCENT),
    ("02", "驱动因素一",   "AI产业链——最强增量引擎",       ACCENT2),
    ("03", "驱动因素二",   "传统下游共振——坚实底盘",       GREEN),
    ("04", "驱动因素三",   "供应链重塑——国产替代加速",     RED),
    ("05", "行业展望",     "景气周期研判与投资建议",         ACCENT),
]

for i, (num, title, sub, color) in enumerate(items):
    y = Inches(0.95 + i * 1.15)
    add_rect(sl, Inches(0.45), y, Inches(12.4), Inches(1.0), fill_color=CARD_BG)
    add_rect(sl, Inches(0.45), y, Inches(0.9),  Inches(1.0), fill_color=color)
    add_textbox(sl, num, Inches(0.45), y + Inches(0.25),
                Inches(0.9), Inches(0.5), font_size=Pt(22), bold=True,
                color=DARK_BG, align=PP_ALIGN.CENTER)
    add_textbox(sl, title, Inches(1.5), y + Inches(0.08),
                Inches(5.5), Inches(0.45), font_size=Pt(18), bold=True, color=WHITE)
    add_textbox(sl, sub, Inches(1.5), y + Inches(0.52),
                Inches(10.5), Inches(0.38), font_size=Pt(13), color=GRAY_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – 市场现状总览
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
section_title(sl, "01  市场现状：三重信号同步亮起")

# Three signal cards
signals = [
    ("📦", "爆  单", ACCENT,
     ["乔锋智能Q1净利润同比 +41.98%",
      "纽威数控Q1营收同比 +25.04%",
      "纽威4月单月订单 >4亿元（历史新高）",
      "日本机床对华订单同比 +57% / +40%",
      "头部产线排产趋满，交期明显延长"]),
    ("⚡", "缺  货", RED,
     ["发那科、三菱数控系统交期大幅延长",
      "部分型号已出现断供情况",
      "丝杠、导轨等传动件同步紧张",
      "AI服务器芯片争夺挤压电子元件产能",
      "上游PCB及伺服电机供应受限"]),
    ("📈", "涨  价", ACCENT2,
     ["铜、铝、塑料等原材料成本上升",
      "供不应求市场格局推动价格走强",
      "日本高端机床率先进入涨价通道",
      "工控产品整机价格上调",
      "行业景气度提升强化乐观预期"]),
]

for i, (icon, label, color, bullets) in enumerate(signals):
    x = Inches(0.45 + i * 4.27)
    card(sl, x, Inches(2.15), Inches(4.05), Inches(4.9), bg=CARD_BG)
    add_rect(sl, x, Inches(2.15), Inches(4.05), Inches(0.6), fill_color=color)
    add_textbox(sl, f"{icon}  {label}",
                x + Inches(0.15), Inches(2.18),
                Inches(3.75), Inches(0.5),
                font_size=Pt(20), bold=True, color=DARK_BG)
    txb = slide_textbox = sl.shapes.add_textbox(
        x + Inches(0.2), Inches(2.85), Inches(3.65), Inches(4.0))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "▸  " + b
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – 爆单数据详解
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
section_title(sl, "01  爆单详解：关键数据一览")

# KPI boxes row 1
kpis = [
    ("+41.98%", "乔锋智能 Q1\n归母净利润同比", ACCENT),
    ("+25.04%", "纽威数控 Q1\n营收同比增速", GREEN),
    (">4亿元", "纽威4月单月订单\n（历史新高）", ACCENT2),
    ("+45%", "纽威4月订单\n同比增速", RED),
]
for i, (val, label, color) in enumerate(kpis):
    x = Inches(0.45 + i * 3.22)
    card(sl, x, Inches(2.1), Inches(3.0), Inches(1.5), bg=CARD_BG)
    add_rect(sl, x, Inches(2.1), Inches(3.0), Inches(0.08), fill_color=color)
    add_textbox(sl, val, x + Inches(0.1), Inches(2.18),
                Inches(2.8), Inches(0.7),
                font_size=Pt(26), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(sl, label, x + Inches(0.1), Inches(2.85),
                Inches(2.8), Inches(0.65),
                font_size=Pt(11), color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# Japan orders section
card(sl, Inches(0.45), Inches(3.8), Inches(5.9), Inches(2.75), bg=CARD_BG)
add_textbox(sl, "🗾  日本机床对华订单（先行指标）",
            Inches(0.65), Inches(3.88), Inches(5.5), Inches(0.45),
            font_size=Pt(14), bold=True, color=ACCENT)

rows = [
    ("2026年1月", "+57%", "同比激增，创近年新高"),
    ("2026年2月", "+40%", "延续强劲增长态势"),
    ("2026年3月", "历史新高", "日本机床整体订单额创历史纪录"),
]
for j, (period, growth, note) in enumerate(rows):
    y = Inches(4.4 + j * 0.65)
    add_rect(sl, Inches(0.55), y, Inches(5.7), Inches(0.55),
             fill_color=RGBColor(0x14, 0x28, 0x40) if j % 2 == 0 else CARD_BG)
    add_textbox(sl, period, Inches(0.65), y + Inches(0.08),
                Inches(1.5), Inches(0.38), font_size=Pt(12), color=GRAY_TEXT)
    add_textbox(sl, growth, Inches(2.2), y + Inches(0.05),
                Inches(1.3), Inches(0.42), font_size=Pt(16), bold=True, color=GREEN)
    add_textbox(sl, note, Inches(3.55), y + Inches(0.1),
                Inches(2.5), Inches(0.35), font_size=Pt(11), color=WHITE)

# Delivery delay section
card(sl, Inches(6.55), Inches(3.8), Inches(6.35), Inches(2.75), bg=CARD_BG)
add_textbox(sl, "🏭  产线排产情况",
            Inches(6.75), Inches(3.88), Inches(6.0), Inches(0.45),
            font_size=Pt(14), bold=True, color=ACCENT2)
txb = sl.shapes.add_textbox(Inches(6.75), Inches(4.38), Inches(5.95), Inches(2.0))
tf = txb.text_frame; tf.word_wrap = True
bullets2 = [
    "头部厂商产线排产趋满，高端机型交付周期延长",
    "广东东莞某企业五轴联动数控机床订单已排至9月份",
    "多家企业反馈在手订单充沛，交货周期大幅拉长",
    "部分企业出现供不应求、无法即时交货状态",
]
first = True
for b in bullets2:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = "▸  " + b
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – 驱动因素一：AI产业链
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT2)
accent_bar(sl, color=ACCENT2)
section_title(sl, "02  驱动因素一：AI产业链——最强增量引擎")

# Two sub-panels
for col, (icon, title, color, points) in enumerate([
    ("🌡️", "AI液冷系统需求爆发", ACCENT, [
        "AI算力爆发使液冷从【可选】变【刚需】",
        "液冷接头（UQD）、液冷板、液冷泵体",
        "要求极高加工精度、密封性与表面粗糙度",
        "直接拉动走心机、数控车床、车铣复合加工中心",
        "高精度磨床采购量大幅增长",
        "2026年Q1液冷订单接近2025全年水平",
    ]),
    ("🤖", "人形机器人产业化加速", GREEN, [
        "产业化进程全面加速，量产节点临近",
        "精密减速器、结构件加工需求爆发",
        "大幅拉动齿轮加工机床需求",
        "五轴联动数控机床成核心设备",
        "特斯拉等头部企业量产带来持续增量",
        "高价值、长周期的确定性需求",
    ]),
]):
    x = Inches(0.45 + col * 6.45)
    card(sl, x, Inches(2.1), Inches(6.2), Inches(4.95), bg=CARD_BG)
    add_rect(sl, x, Inches(2.1), Inches(6.2), Inches(0.65), fill_color=color)
    add_textbox(sl, f"{icon}  {title}",
                x + Inches(0.2), Inches(2.12),
                Inches(5.8), Inches(0.55),
                font_size=Pt(17), bold=True, color=DARK_BG)
    txb = sl.shapes.add_textbox(x + Inches(0.25), Inches(2.85), Inches(5.7), Inches(4.0))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for b in points:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = "▸  " + b
        run.font.size = Pt(12.5)
        run.font.color.rgb = WHITE

# Bottom conclusion bar
add_rect(sl, Inches(0.45), Inches(7.1), Inches(12.4), Inches(0.25),
         fill_color=RGBColor(0x05, 0x3A, 0x5A))
add_textbox(sl, "AI是本轮机床需求爆发的最大增量，已成为不可逆的结构性驱动力",
            Inches(0.6), Inches(7.08), Inches(12.0), Inches(0.3),
            font_size=Pt(12), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – 驱动因素二：传统下游共振
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=GREEN)
accent_bar(sl, color=GREEN)
section_title(sl, "03  驱动因素二：传统下游共振——坚实底盘")

sectors = [
    ("🚢", "船舶制造", GREEN, [
        "Q1新接订单量同比暴增 +195%",
        "手持订单量同比增长 +44%",
        "为大型龙门、重型卧式加工中心",
        "带来中长期确定性需求",
    ]),
    ("✈️", "航空航天", ACCENT, [
        "持续拉动高端五轴加工中心需求",
        "军机、商飞双线景气",
        "精密结构件加工要求极高",
        "国产替代力度持续加大",
    ]),
    ("📱", "消费电子", ACCENT2, [
        "AI手机、折叠屏、AR/VR新品周期",
        "消费电子整体回暖",
        "拉动钻攻机、精雕机需求",
        "小型精密机床迎换机潮",
    ]),
    ("💾", "半导体", RED, [
        "半导体国产替代持续推进",
        "先进制程扩产拉动设备需求",
        "精密加工设备大量引进",
        "本土晶圆厂持续扩建",
    ]),
]

for i, (icon, title, color, points) in enumerate(sectors):
    col = i % 2
    row = i // 2
    x = Inches(0.45 + col * 6.45)
    y = Inches(2.1 + row * 2.55)
    card(sl, x, y, Inches(6.2), Inches(2.35), bg=CARD_BG)
    add_rect(sl, x, y, Inches(0.55), Inches(2.35), fill_color=color)
    add_textbox(sl, icon, x + Inches(0.0), y + Inches(0.75),
                Inches(0.55), Inches(0.7), font_size=Pt(22),
                align=PP_ALIGN.CENTER, color=WHITE)
    add_textbox(sl, title, x + Inches(0.7), y + Inches(0.1),
                Inches(5.3), Inches(0.42), font_size=Pt(15), bold=True, color=color)
    txb = sl.shapes.add_textbox(x + Inches(0.7), y + Inches(0.55), Inches(5.3), Inches(1.65))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for b in points:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "▸  " + b
        run.font.size = Pt(11.5)
        run.font.color.rgb = WHITE

add_rect(sl, Inches(0.45), Inches(7.1), Inches(12.4), Inches(0.25),
         fill_color=RGBColor(0x03, 0x3A, 0x2A))
add_textbox(sl, "多下游行业共振复苏，为机床行业景气度提供坚实底盘支撑",
            Inches(0.6), Inches(7.08), Inches(12.0), Inches(0.3),
            font_size=Pt(12), color=GREEN, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – 驱动因素三：国产替代
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=RED)
accent_bar(sl, color=RED)
section_title(sl, "04  驱动因素三：供应链重塑——国产替代加速")

# Supply shortage analysis
card(sl, Inches(0.45), Inches(2.1), Inches(5.9), Inches(4.95), bg=CARD_BG)
add_rect(sl, Inches(0.45), Inches(2.1), Inches(5.9), Inches(0.6), fill_color=RED)
add_textbox(sl, "⚡  供应链危机分析",
            Inches(0.65), Inches(2.14), Inches(5.5), Inches(0.5),
            font_size=Pt(16), bold=True, color=DARK_BG)

crisis_items = [
    ("触发点", "AI服务器爆发抢占芯片、PCB等电子元件产能"),
    ("传导路径", "电子元件短缺 → 数控系统停产 → 机床交期延长"),
    ("主要品牌", "发那科、三菱等日系品牌交期大幅延长"),
    ("缺口规模", "仅2026年日本系统缺货可释放5-10万台/年缺口"),
    ("外部因素", "欧美系厂商战略收缩，日系供给更加紧张"),
]
txb = sl.shapes.add_textbox(Inches(0.65), Inches(2.82), Inches(5.5), Inches(4.0))
tf = txb.text_frame; tf.word_wrap = True
first = True
for label, content in crisis_items:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = f"【{label}】"
    run.font.size = Pt(11.5)
    run.font.bold = True
    run.font.color.rgb = RED
    run2 = p.add_run()
    run2.text = "  " + content
    run2.font.size = Pt(11.5)
    run2.font.color.rgb = WHITE

# Domestic substitution opportunity
card(sl, Inches(6.55), Inches(2.1), Inches(6.35), Inches(4.95), bg=CARD_BG)
add_rect(sl, Inches(6.55), Inches(2.1), Inches(6.35), Inches(0.6), fill_color=ACCENT)
add_textbox(sl, "🏆  国产替代机遇",
            Inches(6.75), Inches(2.14), Inches(6.0), Inches(0.5),
            font_size=Pt(16), bold=True, color=DARK_BG)

opps = [
    ("替代窗口", "外资断供倒逼客户大规模试用国产系统"),
    ("核心受益", "华中数控（数控系统）、科德数控（五轴机床）"),
    ("不可逆性", "一旦验证可靠性，国产化替代将持续深化"),
    ("长期价值", "中国机床产业加速向价值链中高端攀升"),
    ("政策支持", "【工业母机】战略定位 + 【两新】政策资金支持"),
]
txb = sl.shapes.add_textbox(Inches(6.75), Inches(2.82), Inches(5.95), Inches(4.0))
tf = txb.text_frame; tf.word_wrap = True
first = True
for label, content in opps:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = f"【{label}】"
    run.font.size = Pt(11.5)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run2 = p.add_run()
    run2.text = "  " + content
    run2.font.size = Pt(11.5)
    run2.font.color.rgb = WHITE

add_rect(sl, Inches(0.45), Inches(7.1), Inches(12.4), Inches(0.25),
         fill_color=RGBColor(0x3A, 0x08, 0x14))
add_textbox(sl, "外部压力转化为国内机遇——这是机床国产替代历程中最重要的一次战略窗口",
            Inches(0.6), Inches(7.08), Inches(12.0), Inches(0.3),
            font_size=Pt(12), color=RED, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – 行业展望
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
section_title(sl, "05  行业展望：景气周期研判")

outlook = [
    ("🔄", "周期与结构共振", ACCENT,
     "传统顺周期复苏 + AI新质生产力驱动的结构性行情叠加，持续性与强度有望超越上轮（2020-2023年）新能源周期。"),
    ("📋", "政策强力支持", GREEN,
     "国家【工业母机】战略地位确立，【两新】政策落地为行业提供长期政策保障与资金支持。"),
    ("📅", "景气周期展望", ACCENT2,
     "2026年Q1已成重要业绩拐点，行业有望开启新一轮长达数年的增长周期，当前仍处景气初期。"),
]

for i, (icon, title, color, text) in enumerate(outlook):
    y = Inches(2.1 + i * 1.6)
    card(sl, Inches(0.45), y, Inches(12.4), Inches(1.45), bg=CARD_BG)
    add_rect(sl, Inches(0.45), y, Inches(0.65), Inches(1.45), fill_color=color)
    add_textbox(sl, icon, Inches(0.45), y + Inches(0.35),
                Inches(0.65), Inches(0.65), font_size=Pt(22), align=PP_ALIGN.CENTER, color=WHITE)
    add_textbox(sl, title, Inches(1.25), y + Inches(0.1),
                Inches(4.0), Inches(0.45), font_size=Pt(15), bold=True, color=color)
    add_textbox(sl, text, Inches(1.25), y + Inches(0.55),
                Inches(11.3), Inches(0.8), font_size=Pt(12.5), color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – 投资建议
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT2)
section_title(sl, "05  投资建议：关注三大方向")

categories = [
    ("整机厂商", ACCENT, [
        ("乔锋智能", "液冷、3C、机器人高景气赛道受益"),
        ("津上机床（中国）", "液冷、精密加工领域深度布局"),
        ("纽威数控", "通用复苏，业绩明显拐点，订单创历史新高"),
        ("海天精工", "中大型加工中心，受益制造业全面复苏"),
        ("科德数控", "五轴联动，国产高端标杆，替代空间巨大"),
    ]),
    ("核心系统/部件", GREEN, [
        ("华中数控", "数控系统国产替代龙头，最直接受益标的"),
        ("丝杠、导轨厂商", "核心传动件紧缺，涨价逻辑顺畅"),
        ("伺服电机厂商", "需求旺盛，国产化加速推进"),
    ]),
    ("配套耗材", ACCENT2, [
        ("华锐精密", "硬质合金刀具，随行业景气弹性释放"),
        ("沃尔德", "高性能刀具，技术壁垒显著"),
        ("欧科亿", "刀具销量与利润弹性有望充分释放"),
    ]),
]

for i, (cat_title, color, stocks) in enumerate(categories):
    x = Inches(0.45 + i * 4.27)
    card(sl, x, Inches(2.1), Inches(4.05), Inches(4.95), bg=CARD_BG)
    add_rect(sl, x, Inches(2.1), Inches(4.05), Inches(0.55), fill_color=color)
    add_textbox(sl, cat_title, x + Inches(0.15), Inches(2.12),
                Inches(3.75), Inches(0.48), font_size=Pt(16), bold=True, color=DARK_BG)
    txb = sl.shapes.add_textbox(x + Inches(0.15), Inches(2.75), Inches(3.75), Inches(4.1))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for name, desc in stocks:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = name
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.space_before = Pt(1)
        r2 = p2.add_run()
        r2.text = "    " + desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = GRAY_TEXT

# Disclaimer
add_rect(sl, Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.2),
         fill_color=RGBColor(0x14, 0x1E, 0x30))
add_textbox(sl, "⚠️  以上内容仅供参考，不构成投资建议。投资有风险，入市需谨慎。",
            Inches(0.6), Inches(7.1), Inches(12.0), Inches(0.28),
            font_size=Pt(10), color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – 总结
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg_dark(sl)
add_rect(sl, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT)
add_rect(sl, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), fill_color=ACCENT2)

add_textbox(sl, "总  结",
            Inches(0.45), Inches(0.3), Inches(12.4), Inches(0.7),
            font_size=Pt(28), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

divider_line(sl, Inches(1.05))

# Three pillars
pillars = [
    ("🤖\nAI算力\n基础设施建设",    ACCENT,
     "液冷系统与人形机器人催生持续高价值增量需求，替代传统新能源成为最大驱动力"),
    ("🏭\n制造业\n全面复苏",         GREEN,
     "船舶、航空、消费电子、半导体多行业共振，传统基本盘全面筑底反转"),
    ("🇨🇳\n供应链\n自主可控",         ACCENT2,
     "外资断供危机转化为国产替代战略机遇，不可逆的国产化进程正式提速"),
]

for i, (title, color, desc) in enumerate(pillars):
    x = Inches(0.55 + i * 4.27)
    card(sl, x, Inches(1.3), Inches(4.0), Inches(3.8), bg=CARD_BG)
    add_rect(sl, x, Inches(1.3), Inches(4.0), Inches(0.06), fill_color=color)
    add_textbox(sl, title, x + Inches(0.1), Inches(1.42),
                Inches(3.8), Inches(1.4),
                font_size=Pt(16), bold=True, color=color,
                align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + Inches(0.15), Inches(2.85),
                Inches(3.7), Inches(2.0),
                font_size=Pt(12), color=WHITE, wrap=True, align=PP_ALIGN.CENTER)

# Final statement
add_rect(sl, Inches(0.45), Inches(5.25), Inches(12.4), Inches(1.55),
         fill_color=RGBColor(0x05, 0x3A, 0x5A), line_color=ACCENT, line_width=Pt(1.5))
add_textbox(sl, "【核心结论】",
            Inches(0.65), Inches(5.32), Inches(11.8), Inches(0.38),
            font_size=Pt(14), bold=True, color=ACCENT)
add_textbox(sl,
    "当前机床行业的火热并非短期脉冲，而是 AI算力基础设施建设、制造业全面复苏、"
    "供应链自主可控 三大趋势共同推动的结构性机遇。\n"
    "2026年一季度已成重要业绩拐点，行业有望开启新一轮长达数年的增长周期。",
    Inches(0.65), Inches(5.72), Inches(11.8), Inches(1.0),
    font_size=Pt(13.5), color=WHITE, wrap=True)

add_textbox(sl, "报告日期：2026年5月13日  ·  数据来源：公司公告、行业研究、媒体报道",
            Inches(0.45), Inches(6.95), Inches(12.4), Inches(0.35),
            font_size=Pt(10), color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
OUTPUT = "/home/user/UnoxyRich/机床市场深度分析_2026Q1.pptx"
prs.save(OUTPUT)
print(f"✅  PPT saved: {OUTPUT}")
